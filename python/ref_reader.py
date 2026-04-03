"""참고자료 텍스트 추출기.
지원: .txt, .csv, .xlsx, .json, .md, .pdf, .html, .xml
추가: .docx, .pptx, .doc, .ppt, .rtf 등 → PDF 변환 후 텍스트 추출
HWP/HWPX는 hwp_analyzer.analyze_document 사용 (이 모듈에서는 다루지 않음)
"""
import os
import sys
import json
import subprocess
import tempfile


def read_reference(file_path, max_chars=30000):
    """참고자료 파일에서 텍스트 추출."""
    file_path = os.path.abspath(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext in ('.txt', '.md', '.log'):
        return _read_text(file_path, max_chars)
    elif ext == '.csv':
        return _read_csv(file_path, max_chars)
    elif ext in ('.xlsx', '.xls'):
        return _read_excel(file_path, max_chars)
    elif ext == '.json':
        return _read_json(file_path, max_chars)
    elif ext == '.pdf':
        return _read_pdf(file_path, max_chars)
    elif ext in ('.html', '.htm'):
        return _read_html(file_path, max_chars)
    elif ext == '.xml':
        return _read_xml(file_path, max_chars)
    elif ext in ('.docx', '.doc', '.pptx', '.ppt', '.rtf', '.odt', '.odp'):
        return _read_via_pdf_conversion(file_path, max_chars)
    else:
        raise ValueError(
            f"지원하지 않는 파일 형식: {ext}. "
            f"지원: .txt, .md, .csv, .xlsx, .json, .pdf, .html, .xml, .docx, .pptx"
        )


def _read_text(path, max_chars):
    with open(path, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read(max_chars)
    return {
        "format": "text",
        "file_name": os.path.basename(path),
        "content": content,
        "char_count": len(content),
    }


def _read_csv(path, max_chars):
    import csv
    rows = []
    total_chars = 0
    with open(path, 'r', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f)
        for row in reader:
            row_text = ','.join(row)
            total_chars += len(row_text)
            if total_chars > max_chars:
                break
            rows.append(row)

    headers = rows[0] if rows else []
    data = rows[1:] if len(rows) > 1 else []
    return {
        "format": "csv",
        "file_name": os.path.basename(path),
        "headers": headers,
        "data": data,
        "row_count": len(data),
    }


def _read_excel(path, max_chars):
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl이 필요합니다. pip install openpyxl")

    wb = None
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheets = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            total_chars = 0
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                total_chars += sum(len(c) for c in row_data)
                if total_chars > max_chars:
                    break
                rows.append(row_data)

            headers = rows[0] if rows else []
            data = rows[1:] if len(rows) > 1 else []
            sheets.append({
                "sheet_name": sheet_name,
                "headers": headers,
                "data": data,
                "row_count": len(data),
            })

        return {
            "format": "excel",
            "file_name": os.path.basename(path),
            "sheets": sheets,
            "sheet_count": len(sheets),
        }
    finally:
        if wb:
            wb.close()


def _read_json(path, max_chars):
    with open(path, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read(max_chars)

    data = json.loads(content)
    return {
        "format": "json",
        "file_name": os.path.basename(path),
        "data": data,
    }


def _read_pdf(path, max_chars):
    """PDF에서 텍스트 추출 (PyMuPDF 사용)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("PyMuPDF가 필요합니다. pip install PyMuPDF")

    doc = fitz.open(path)
    pages = []
    total_chars = 0
    for i, page in enumerate(doc):
        text = page.get_text("text")
        total_chars += len(text)
        pages.append({"page": i + 1, "text": text})
        if total_chars > max_chars:
            break
    doc.close()

    full_text = "\n\n".join(p["text"] for p in pages)
    return {
        "format": "pdf",
        "file_name": os.path.basename(path),
        "content": full_text[:max_chars],
        "page_count": len(pages),
        "char_count": len(full_text[:max_chars]),
    }


def _read_via_pdf_conversion(path, max_chars):
    """DOCX/PPTX 등 비지원 확장자 → PDF 변환 후 텍스트 추출."""
    ext = os.path.splitext(path)[1].lower()

    # 1순위: LibreOffice CLI로 PDF 변환
    pdf_path = _convert_to_pdf_libreoffice(path)
    if pdf_path:
        result = _read_pdf(pdf_path, max_chars)
        result["original_format"] = ext.lstrip('.')
        result["conversion_method"] = "libreoffice"
        # 임시 PDF 삭제
        try:
            os.remove(pdf_path)
        except Exception:
            pass
        return result

    # 2순위: python-docx로 직접 텍스트 추출 (DOCX만)
    if ext == '.docx':
        result = _read_docx_direct(path, max_chars)
        if result:
            return result

    # 3순위: python-pptx로 직접 텍스트 추출 (PPTX만)
    if ext == '.pptx':
        result = _read_pptx_direct(path, max_chars)
        if result:
            return result

    raise ValueError(
        f"{ext} 파일을 읽을 수 없습니다. "
        f"LibreOffice를 설치하면 자동 변환됩니다: https://www.libreoffice.org/download/"
    )


def _convert_to_pdf_libreoffice(path):
    """LibreOffice CLI로 PDF 변환. 성공 시 PDF 경로 반환, 실패 시 None."""
    # LibreOffice 경로 탐색
    soffice_paths = [
        "soffice",  # PATH에 있으면
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]

    soffice = None
    for p in soffice_paths:
        try:
            subprocess.run([p, "--version"], capture_output=True, timeout=5)
            soffice = p
            break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    if not soffice:
        print("[INFO] LibreOffice 미설치 — PDF 변환 불가, 대체 방법 시도", file=sys.stderr)
        return None

    try:
        outdir = tempfile.gettempdir()
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", outdir, path],
            capture_output=True, timeout=60
        )
        basename = os.path.splitext(os.path.basename(path))[0]
        pdf_path = os.path.join(outdir, f"{basename}.pdf")
        if os.path.exists(pdf_path):
            return pdf_path
    except Exception as e:
        print(f"[WARN] LibreOffice 변환 실패: {e}", file=sys.stderr)

    return None


def _read_docx_direct(path, max_chars):
    """python-docx로 DOCX 텍스트 직접 추출."""
    try:
        from docx import Document
    except ImportError:
        return None

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    content = "\n".join(paragraphs)[:max_chars]
    return {
        "format": "docx",
        "file_name": os.path.basename(path),
        "content": content,
        "paragraph_count": len(paragraphs),
        "char_count": len(content),
    }


def _read_pptx_direct(path, max_chars):
    """python-pptx로 PPTX 텍스트 직접 추출."""
    try:
        from pptx import Presentation
    except ImportError:
        return None

    prs = Presentation(path)
    slides = []
    total_chars = 0
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        slide_text = "\n".join(texts)
        total_chars += len(slide_text)
        slides.append({"slide": i + 1, "text": slide_text})
        if total_chars > max_chars:
            break

    full_text = "\n\n".join(s["text"] for s in slides)
    return {
        "format": "pptx",
        "file_name": os.path.basename(path),
        "content": full_text[:max_chars],
        "slide_count": len(slides),
        "char_count": len(full_text[:max_chars]),
    }


def _read_html(path, max_chars):
    """HTML 파일에서 텍스트 + 표 추출."""
    import re
    with open(path, 'r', encoding='utf-8', errors='replace') as f:
        html = f.read(max_chars * 3)  # HTML 태그 포함이므로 더 많이 읽음

    # 표 추출 (정규식 기반 — BeautifulSoup 없이)
    tables = []
    table_pattern = re.compile(r'<table[^>]*>(.*?)</table>', re.DOTALL | re.IGNORECASE)
    for match in table_pattern.finditer(html):
        table_html = match.group(1)
        rows = []
        row_pattern = re.compile(r'<tr[^>]*>(.*?)</tr>', re.DOTALL | re.IGNORECASE)
        for row_match in row_pattern.finditer(table_html):
            row_html = row_match.group(1)
            cells = re.findall(r'<t[dh][^>]*>(.*?)</t[dh]>', row_html, re.DOTALL | re.IGNORECASE)
            # 태그 제거
            clean_cells = [re.sub(r'<[^>]+>', '', c).strip() for c in cells]
            if clean_cells:
                rows.append(clean_cells)
        if rows:
            headers = rows[0] if rows else []
            data = rows[1:] if len(rows) > 1 else []
            tables.append({"headers": headers, "data": data, "row_count": len(data)})

    # 본문 텍스트 (태그 제거)
    text = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()[:max_chars]

    return {
        "format": "html",
        "file_name": os.path.basename(path),
        "content": text,
        "tables": tables,
        "table_count": len(tables),
        "char_count": len(text),
    }


def _read_xml(path, max_chars):
    """XML 파일에서 구조화된 데이터 추출."""
    import xml.etree.ElementTree as ET

    tree = ET.parse(path)
    root = tree.getroot()

    # 네임스페이스 제거 (간편 접근)
    def strip_ns(tag):
        return tag.split('}')[-1] if '}' in tag else tag

    def elem_to_dict(elem):
        """XML 요소를 dict로 변환 (재귀)."""
        result = {}
        # 속성
        if elem.attrib:
            result["@attributes"] = dict(elem.attrib)
        # 텍스트
        if elem.text and elem.text.strip():
            result["@text"] = elem.text.strip()
        # 자식 요소
        children = {}
        for child in elem:
            tag = strip_ns(child.tag)
            child_data = elem_to_dict(child)
            if tag in children:
                # 같은 태그 반복 → 리스트화
                if not isinstance(children[tag], list):
                    children[tag] = [children[tag]]
                children[tag].append(child_data)
            else:
                children[tag] = child_data
        result.update(children)
        return result

    data = {strip_ns(root.tag): elem_to_dict(root)}

    # 텍스트 미리보기
    text_parts = []
    for elem in root.iter():
        if elem.text and elem.text.strip():
            text_parts.append(elem.text.strip())
    content = "\n".join(text_parts)[:max_chars]

    return {
        "format": "xml",
        "file_name": os.path.basename(path),
        "data": data,
        "content": content,
        "char_count": len(content),
    }
