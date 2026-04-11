"""ref_reader.conversion — 변환/fallback 리더.

1차 지원 안 되는 포맷 (.docx, .pptx, .rtf, .odt, .odp 등) 을 처리:
- LibreOffice CLI 로 PDF 변환 후 _read_pdf 재사용 (1순위)
- python-docx 직접 추출 (2순위, DOCX 만)
- python-pptx 직접 추출 (3순위, PPTX 만)

공개 API:
- _read_via_pdf_conversion — 변환 파이프라인 진입점
"""
import os
import sys
import subprocess
import tempfile

from .readers import _read_pdf


def _read_via_pdf_conversion(path, max_chars):
    """DOCX/PPTX 등 비지원 확장자 → PDF 변환 후 텍스트 추출."""
    ext = os.path.splitext(path)[1].lower()

    # 1순위: LibreOffice CLI로 PDF 변환
    pdf_path = _convert_to_pdf_libreoffice(path)
    if pdf_path:
        result = _read_pdf(pdf_path, max_chars)
        result["original_format"] = ext.lstrip('.')
        result["conversion_method"] = "libreoffice"
        # 임시 PDF 삭제
        try:
            os.remove(pdf_path)
        except Exception:
            pass
        return result

    # 2순위: python-docx로 직접 텍스트 추출 (DOCX만)
    if ext == '.docx':
        result = _read_docx_direct(path, max_chars)
        if result:
            return result

    # 3순위: python-pptx로 직접 텍스트 추출 (PPTX만)
    if ext == '.pptx':
        result = _read_pptx_direct(path, max_chars)
        if result:
            return result

    raise ValueError(
        f"{ext} 파일을 읽을 수 없습니다. "
        f"LibreOffice를 설치하면 자동 변환됩니다: https://www.libreoffice.org/download/"
    )


def _convert_to_pdf_libreoffice(path):
    """LibreOffice CLI로 PDF 변환. 성공 시 PDF 경로 반환, 실패 시 None."""
    # LibreOffice 경로 탐색
    soffice_paths = [
        "soffice",  # PATH에 있으면
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]

    soffice = None
    for p in soffice_paths:
        try:
            subprocess.run([p, "--version"], capture_output=True, timeout=5)
            soffice = p
            break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    if not soffice:
        print("[INFO] LibreOffice 미설치 — PDF 변환 불가, 대체 방법 시도", file=sys.stderr)
        return None

    try:
        outdir = tempfile.gettempdir()
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", outdir, path],
            capture_output=True, timeout=60
        )
        basename = os.path.splitext(os.path.basename(path))[0]
        pdf_path = os.path.join(outdir, f"{basename}.pdf")
        if os.path.exists(pdf_path):
            return pdf_path
    except Exception as e:
        print(f"[WARN] LibreOffice 변환 실패: {e}", file=sys.stderr)

    return None


def _read_docx_direct(path, max_chars):
    """python-docx로 DOCX 텍스트 직접 추출."""
    try:
        from docx import Document
    except ImportError:
        return None

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    content = "\n".join(paragraphs)[:max_chars]
    return {
        "format": "docx",
        "file_name": os.path.basename(path),
        "content": content,
        "paragraph_count": len(paragraphs),
        "char_count": len(content),
    }


def _read_pptx_direct(path, max_chars):
    """python-pptx로 PPTX 텍스트 직접 추출."""
    try:
        from pptx import Presentation
    except ImportError:
        return None

    prs = Presentation(path)
    slides = []
    total_chars = 0
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        slide_text = "\n".join(texts)
        total_chars += len(slide_text)
        slides.append({"slide": i + 1, "text": slide_text})
        if total_chars > max_chars:
            break

    full_text = "\n\n".join(s["text"] for s in slides)
    return {
        "format": "pptx",
        "file_name": os.path.basename(path),
        "content": full_text[:max_chars],
        "slide_count": len(slides),
        "char_count": len(full_text[:max_chars]),
    }
